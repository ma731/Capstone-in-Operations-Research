"""Build the capstone defense deck (16:9 PowerPoint) from the IE design system.

Generates ``deck/capstone_defense.pptx``: an on-brand, editable ~18-slide talk that
reuses the thesis figures and the poster palette. Equations are rendered crisply with
matplotlib mathtext into ``deck/_eq_*.png`` so they stay sharp at any zoom.

Run:  .venv/Scripts/python scripts/build_deck.py
"""
from __future__ import annotations

import os
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
DECK = ROOT / "deck"
DECK.mkdir(exist_ok=True)

# ---- brand palette (matches poster/design-system.css) ----
NAVY = RGBColor(0x0E, 0x2A, 0x52)
NAVY2 = RGBColor(0x1F, 0x3B, 0x63)
GOLD = RGBColor(0xE6, 0x9F, 0x00)
RUST = RGBColor(0xB3, 0x40, 0x2F)
SAGE = RGBColor(0x4A, 0x7C, 0x59)
INK = RGBColor(0x16, 0x20, 0x2E)
MUTED = RGBColor(0x5B, 0x66, 0x75)
LINE = RGBColor(0xD9, 0xDE, 0xE6)
TINT = RGBColor(0xF3, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD = RGBColor(0xCD, 0xD6, 0xE4)

SERIF = "Georgia"          # display / titles (universally present on Windows)
SANS = "Segoe UI"          # body / UI

EMU_IN = 914400
SW, SH = 13.333, 7.5       # 16:9 widescreen inches

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------- helpers
def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _oval(slide, x, y, d, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, line_spacing=1.06):
    """runs: list of paragraphs; each paragraph is a list of (text, font, size, color, bold)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, font, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def _bullets(slide, x, y, w, h, items, *, size=18, gap=10, color=INK, lead=None):
    """items: list of (text, color_or_None, bold). Renders gold-dot bullets."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    first = True
    if lead:
        p = tf.paragraphs[0]
        first = False
        p.space_after = Pt(gap + 2)
        p.line_spacing = 1.12
        r = p.add_run(); r.text = lead
        r.font.name = SANS; r.font.size = Pt(size + 3); r.font.bold = True; r.font.color.rgb = NAVY
    for (txt, col, bold) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = 1.12
        rd = p.add_run(); rd.text = "●  "
        rd.font.name = SANS; rd.font.size = Pt(size - 4); rd.font.color.rgb = GOLD; rd.font.bold = True
        rt = p.add_run(); rt.text = txt
        rt.font.name = SANS; rt.font.size = Pt(size)
        rt.font.color.rgb = col or color; rt.font.bold = bold
    return tb


def _pic(slide, path, x, y, *, w=None, h=None, max_w=None, max_h=None, center_x=None):
    """Place an image preserving aspect. Provide w or h, or a max box (max_w,max_h)."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    if max_w and max_h:
        if max_w / max_h > ar:
            h = max_h; w = max_h * ar
        else:
            w = max_w; h = max_w / ar
    elif w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    if center_x is not None:
        x = center_x - w / 2
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def _title_bar(slide, title, num):
    _rect(slide, 0, 0, SW, 1.04, NAVY)
    _rect(slide, 0, 1.04, SW, 0.06, GOLD)
    # number badge
    _oval(slide, 0.5, 0.27, 0.5, GOLD)
    _text(slide, 0.5, 0.27, 0.5, 0.5, [[(str(num), SERIF, 22, NAVY, True)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, 1.2, 0.0, SW - 2.4, 1.04, [[(title, SERIF, 27, WHITE, True)]],
          anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, num):
    _rect(slide, 0.5, SH - 0.42, SW - 1.0, 0.012, LINE)
    _text(slide, 0.5, SH - 0.40, 8.0, 0.34,
          [[("The Price of Sophistication: carbon-aware data-center scheduling", SANS, 10, MUTED, False)]],
          anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, SW - 2.5, SH - 0.40, 2.0, 0.34,
          [[("M. Ortiz Togashi  ·  " + str(num), SANS, 10, MUTED, False)]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def content(title, num):
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, WHITE)
    _title_bar(s, title, num)
    _footer(s, num)
    return s


# ----------------------------------------------------------------------------- equations
def render_eq(tex, name, *, size=30, color="#0E2A52"):
    out = DECK / f"_eq_{name}.png"
    fig = plt.figure(figsize=(0.1, 0.1))
    fig.text(0.0, 0.0, f"${tex}$", fontsize=size, color=color)
    fig.savefig(out, dpi=300, transparent=True, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    return out


EQ_DRO = render_eq(
    r"\min_{x \in \mathcal{X}}\ \langle \bar{\rho},\, x \rangle \;+\; "
    r"\varepsilon\, \| L^{\top} x \|_2 \qquad LL^{\top} = \hat{\Sigma}", "dro")

# =============================================================================== SLIDES

# ---- 1. Title ----
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_rect(s, 0, 0, 0.28, SH, GOLD)               # gold spine
# white IE logo if available
logo = ROOT / "poster" / "figs" / "ie_logo_white.png"
if logo.exists():
    _pic(s, logo, 0.95, 0.72, w=2.0)
_text(s, 0.95, 2.42, 11.7, 2.4, [
    [("The Price of Sophistication", SERIF, 33, WHITE, True)],
    [("When Do Spatial and Robust Models Pay in", SERIF, 25, WHITE, True)],
    [("Carbon-Aware Data-Center Scheduling?", SERIF, 25, WHITE, True)],
], line_spacing=1.04, space_after=8)
_rect(s, 1.0, 4.92, 3.2, 0.05, GOLD)
_text(s, 0.95, 5.12, 11.4, 0.8,
      [[("A working day-ahead migration scheduler, and a decision rule for when each "
         "modelling layer earns its complexity", SANS, 17, CLOUD, False)]])
_text(s, 0.95, 6.15, 11.0, 1.1, [
    [("Marco Ortiz Togashi", SANS, 17, WHITE, True),
     ("    ·    IE University  ·  MSc Business Analytics & Data Science", SANS, 15, CLOUD, False)],
    [("Supervisor: Prof. Bissan Ghaddar", SANS, 15, CLOUD, False)],
], space_after=4)

# ---- 2. The opportunity ----
s = content("Flexible compute meets a varying grid", 2)
_bullets(s, 0.7, 1.5, 6.0, 5.0, [
    ("Compute is unusually flexible: batch and AI-training jobs can shift in time "
     "and across sites.", None, False),
    ("Carbon intensity of electricity varies sharply by hour and by region, driven by "
     "wind, solar, hydro, and demand.", None, False),
    ("So an operator can ship work to cleaner hours and cleaner regions, cutting "
     "emissions at almost no cost.", None, False),
    ("The open question is not whether to schedule, but which modelling layers actually "
     "pay for their complexity.", NAVY, True),
], size=19, gap=16, lead="The opportunity")
_pic(s, FIG / "correlation_map.png", 7.0, 1.7, max_w=5.7, max_h=4.6)
_text(s, 7.0, 6.35, 5.7, 0.4, [[("Carbon intensity is spatially structured, and it moves.",
      SANS, 12, MUTED, False)]], align=PP_ALIGN.CENTER)

# ---- 3. The question: which layers pay ----
s = content("Which layers earn their complexity?", 3)
_bullets(s, 0.7, 1.55, 11.9, 3.6, [
    ("Each layer (inter-region transfer, spatial covariance, distributional "
     "robustness, richer copulas) adds estimation and compute cost.", None, False),
    ("The literature has assumed this sophistication pays, rather than measuring it.", None, False),
    ("We build a working day-ahead scheduler and price each layer directly, "
     "out of sample.", NAVY, True),
], size=20, gap=16, lead="Sophistication is not free")
_rect(s, 0.7, 5.45, 11.9, 1.25, TINT)
_text(s, 1.0, 5.52, 11.3, 1.1, [[(
    "RQ1: how much does the transfer lever save?   RQ2: does passive covariance or a "
    "copula add anything?   RQ3: when does robustness pay?", SERIF, 18, NAVY, True)]],
    anchor=MSO_ANCHOR.MIDDLE)

# ---- 4. The day-ahead system ----
s = content("The day-ahead scheduler", 4)
_bullets(s, 0.7, 1.5, 6.1, 5.0, [
    ("Forecast tomorrow's per-region carbon intensity from history.", None, False),
    ("Optimize a migration schedule: move splittable compute across regions under "
     "ramp limits and a transfer budget Phi.", None, False),
    ("Roll forward day by day; score out-of-sample CVaR(0.95) of daily emissions.", None, False),
    ("Honest baseline: carbon-aware scheduling with no inter-region transfer (Phi=0), "
     "on the same feasible set.", NAVY, True),
], size=18, gap=14, lead="Forecast, optimize, roll")
_pic(s, FIG / "model_validation.png", 7.0, 1.6, max_w=5.7, max_h=4.6)
_text(s, 7.0, 6.3, 5.7, 0.4, [[(
    "The scheduler matches the exact optimum on small instances.", SANS, 12, MUTED, False)]],
    align=PP_ALIGN.CENTER)

# ---- 5. The value: transfer lever (RQ1) ----
s = content("RQ1: the transfer lever pays", 5)
_pic(s, FIG / "transfer_value_curve.png", 0.7, 1.55, max_w=7.6, max_h=4.9)
_bullets(s, 8.5, 1.7, 4.2, 4.6, [
    ("Active migration cuts out-of-sample CVaR(0.95) by 4.0 to 9.9% over Phi=0.", SAGE, True),
    ("Per grid: Western 4.0%, Eastern 9.9%, Diversified 9.0%.", None, False),
    ("The lever is the spatial mean: at each hour one region is cleanest, "
     "and migration ships work there.", None, False),
    ("Deterministic: captured with no dependence model at all.", NAVY, True),
], size=17, gap=14, lead="The lever, and it works")

# ---- 6. The complexity-value frontier ----
s = content("Where the value lives", 6)
_pic(s, FIG / "complexity_frontier.png", 0.7, 1.55, max_w=7.6, max_h=4.9)
_bullets(s, 8.5, 1.7, 4.2, 4.6, [
    ("Every model class on one complexity-value frontier.", None, False),
    ("Deterministic transfer dominates the frontier.", SAGE, True),
    ("Passive covariance and copula layers add no height.", None, False),
    ("Robust transfer pays only in the tail, conditional on severity.", None, False),
], size=18, gap=15, lead="The frontier")

# ---- 7. The screening rule (RQ2) ----
s = content("RQ2: passive sophistication adds nothing", 7)
_pic(s, FIG / "cv_curve.png", 0.7, 1.55, max_w=7.4, max_h=4.9)
_bullets(s, 8.4, 1.7, 4.3, 4.8, [
    ("The spatial covariance gap stays below 0.4% of CVaR, and its sign is not "
     "stable.", None, False),
    ("The DRO is genuinely engaged (CV selects a non-trivial radius), so this is an "
     "active null, not a switched-off one.", None, False),
    ("Design justification: a simple per-region model is enough.", NAVY, True),
], size=17, gap=14, lead="A strong, active null")

# ---- 8. Why the screening rule holds: mechanism ----
s = content("Why covariance cannot add value", 8)
_bullets(s, 0.7, 1.5, 6.1, 4.6, [
    ("The mean-leveling test: flatten the mean field and the covariance suddenly pays, up to "
     "+1.46%.", RUST, True),
    ("So the cross-region signal is real, but masked by the within-day mean carbon field.", None, False),
    ("Residual dependence is also non-elliptical and upper-tail independent, "
     "invisible to a covariance ball.", None, False),
    ("We checked richer dependence models (Gaussian, Clayton, comonotone): they add "
     "nothing.", None, False),
], size=18, gap=13, lead="The mechanism")
_rect(s, 7.0, 1.7, 5.7, 4.5, TINT)
_text(s, 7.25, 1.95, 5.2, 4.1, [
    [("Mean-dominance bound", SERIF, 19, NAVY, True)],
    [("Carbon intensity is driven overwhelmingly by the time-of-day mean (solar, "
      "demand): shared, predictable, and large.", SANS, 16, INK, False)],
    [("The cross-region covariance is a second-order wrinkle on a first-order wave, "
      "so it cannot move the optimal schedule. CVaR's translation invariance makes "
      "the dependence model second-order to the mean.", SANS, 16, INK, False)],
], space_after=12)

# ---- 9. Why DRO: day-ahead forecast error ----
s = content("RQ3: why robustify at all", 9)
_bullets(s, 0.7, 1.5, 6.1, 4.6, [
    ("Scheduling against tomorrow's carbon means scheduling under forecast error.", None, False),
    ("We hedge that error with a data-driven Wasserstein DRO: a daily, rolling "
     "ambiguity set around the empirical distribution.", None, False),
    ("It reduces to a second-order cone program, solved exactly and fast.", None, False),
    ("This is a price-of-robustness question: it buys tail reduction at a mean "
     "premium.", NAVY, True),
], size=18, gap=14, lead="Hedging day-ahead forecast error")
_rect(s, 7.0, 2.3, 5.7, 1.5, TINT)
_pic(s, EQ_DRO, 7.0, 2.65, max_w=5.4, max_h=0.95, center_x=9.85)
_text(s, 7.0, 4.1, 5.7, 2.0, [[(
    "The mean term carries the saving; the cone term "
    "epsilon times the norm of L-transpose x prices the hedge. "
    "The value of the stochastic solution is near zero under nominal carbon.",
    SANS, 16, MUTED, False)]])

# ---- 10. The price-of-robustness crossover ----
s = content("The price-of-robustness crossover", 10)
_pic(s, FIG / "crossover.png", 0.7, 1.55, max_w=7.6, max_h=4.9)
_bullets(s, 8.5, 1.7, 4.2, 4.8, [
    ("Robust beats risk-neutral only past an emergency severity M* near 3.", None, False),
    ("Real grids, tested across 17 zones, peak at only M near 1.4.", RUST, True),
    ("So robustness does not activate on observed data; even Winter Storm Uri "
     "reached only 1.3x.", None, False),
    ("On real data, deterministic transfer is unambiguously dominant.", NAVY, True),
], size=17, gap=13, lead="M* near 3, real grids stay below")

# ---- 11. The decision rule ----
s = content("The decision rule", 11)
_rect(s, 0.7, 1.5, 5.85, 2.45, TINT)
_rect(s, 0.7, 1.5, 0.14, 2.45, SAGE)
_text(s, 1.0, 1.66, 5.4, 0.5, [[("Layer 1: always", SERIF, 19, NAVY, True)]])
_text(s, 1.0, 2.18, 5.4, 1.6, [[(
    "Schedule per-region, day-ahead. Temporal shifting captures the diurnal "
    "carbon cycle.", SANS, 16, INK, False)]])
_rect(s, 6.75, 1.5, 5.85, 2.45, TINT)
_rect(s, 6.75, 1.5, 0.14, 2.45, GOLD)
_text(s, 7.05, 1.66, 5.4, 0.5, [[("Layer 2: the lever", SERIF, 19, NAVY, True)]])
_text(s, 7.05, 2.18, 5.4, 1.6, [[(
    "Add inter-region transfer if migration bandwidth exists: a 4.0 to 9.9% CVaR "
    "reduction over Phi=0. This is where the value is.", SANS, 16, INK, False)]])
_rect(s, 0.7, 4.15, 5.85, 2.45, TINT)
_rect(s, 0.7, 4.15, 0.14, 2.45, RUST)
_text(s, 1.0, 4.31, 5.4, 0.5, [[("Layer 3: skip", SERIF, 19, NAVY, True)]])
_text(s, 1.0, 4.83, 5.4, 1.6, [[(
    "Spatial covariance and copulas. They add below 0.4% of CVaR: not worth the "
    "complexity.", SANS, 16, INK, False)]])
_rect(s, 6.75, 4.15, 5.85, 2.45, NAVY)
_rect(s, 6.75, 4.15, 0.14, 2.45, GOLD)
_text(s, 7.05, 4.31, 5.4, 0.5, [[("Layer 4: conditional", SERIF, 19, WHITE, True)]])
_text(s, 7.05, 4.83, 5.4, 1.6, [[(
    "Robustify only if you expect emergencies past M* near 3. Real grids do not "
    "reach it.", SANS, 16, CLOUD, False)]])

# ---- 12. Data ----
s = content("Data: three grids across the spectrum", 12)
_bullets(s, 0.7, 1.5, 6.1, 4.8, [
    ("US West: a strongly, uniformly correlated grid (common-mode).", None, False),
    ("Eastern Interconnection belt: an Ontario-anchored mid-correlation set.", None, False),
    ("Engineered solar / wind / hydro portfolio: low, heterogeneous correlation.", None, False),
    ("Iberia-France: an independent low-correlation external-validity anchor.", None, False),
], size=18, gap=14, lead="Real Electricity Maps carbon intensity")
_pic(s, FIG / "ci_corr_heatmap_us_west.png", 7.0, 1.6, max_w=5.7, max_h=4.6)
_text(s, 7.0, 6.3, 5.7, 0.4, [[(
    "Correlation is real, up to 0.78 on the US West.", SANS, 12, MUTED, False)]],
    align=PP_ALIGN.CENTER)

# ---- 13. Validation / test suite ----
s = content("Validation and reproducibility", 13)
_bullets(s, 0.7, 1.5, 6.1, 4.9, [
    ("The null survives Ledoit-Wolf shrinkage, residualization, and "
     "Benjamini-Hochberg correction across cells.", None, False),
    ("Walk-forward out-of-sample validation and a per-cell equivalence test.", None, False),
    ("199 unit tests (186 in CI); bootstrap confidence intervals on every gap.", NAVY, True),
    ("Pre-registered, version-controlled, with archived summary tables for every "
     "reported number.", None, False),
], size=18, gap=14, lead="A pre-registered battery")
_pic(s, FIG / "robustness.png", 7.0, 1.7, max_w=5.7, max_h=4.6)

# ---- 14. Limitations ----
s = content("Limitations, stated honestly", 14)
_bullets(s, 0.7, 1.5, 11.9, 4.9, [
    ("One canonical workload (splittable load, ramp limits); a sharply different cost "
     "geometry could re-weight the covariance term.", None, False),
    ("Primary read is a single held-out year; walk-forward reproduces it, but "
     "multi-year evaluation would strengthen external validity.", None, False),
    ("Results generalize to grids where the diurnal mean dominates residual "
     "covariance: an empirical regularity, not a universal law.", None, False),
    ("The crossover is characterized, not observed: a grid with severe emergencies "
     "past M* near 3 would activate the robust layer.", None, False),
], size=19, gap=15, lead="Where the claims could be pushed")

# ---- 15. Conclusions ----
s = content("Conclusions and contributions", 15)
_bullets(s, 0.7, 1.5, 11.9, 4.9, [
    ("A working day-ahead migration scheduler with honest savings: a 4.0 to 9.9% CVaR "
     "reduction over the Phi=0 baseline, the dominant lever.", SAGE, True),
    ("A screening rule: passive covariance and copulas add nothing, with a "
     "mean-dominance bound that says why.", None, False),
    ("A price-of-robustness decision rule with a tested, data-grounded bound "
     "(M near 1.4 below M* near 3).", None, False),
    ("A reproducible, pre-registered pipeline under 199 unit tests.", NAVY, True),
], size=19, gap=15, lead="What this thesis delivers")

# ---- 16. The decision rule recap + thanks ----
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_rect(s, 0, 0, 0.28, SH, GOLD)
_text(s, 0.95, 0.9, 11.4, 1.0, [[("The takeaway, and thank you", SERIF, 32, WHITE, True)]])
_rect(s, 0.95, 1.78, 2.6, 0.05, GOLD)
_text(s, 0.95, 2.15, 11.4, 0.5, [[("The rule in one line", SANS, 19, GOLD, True)]])
_bullets(s, 0.95, 2.75, 11.4, 3.0, [
    ("Always schedule per-region.", CLOUD, False),
    ("Add inter-region transfer: this is the lever, 4.0 to 9.9% over Phi=0.", CLOUD, False),
    ("Skip covariance and copulas: they add nothing.", CLOUD, False),
    ("Robustify only if you expect emergencies past M* near 3.", CLOUD, False),
], size=18, gap=14, color=CLOUD)
_text(s, 0.95, 6.4, 11.4, 0.8, [[(
    "Thank you. Questions welcome.", SERIF, 22, GOLD, True)]])

# ---- 17. Backup: the falsification test ----
s = content("Backup: the falsification test", 17)
_bullets(s, 0.7, 1.45, 6.1, 5.0, [
    ("Fit the schedule to the full joint covariance.", None, False),
    ("Refit to a block-diagonal covariance: identical marginals, cross-region "
     "structure destroyed (shuffled).", None, False),
    ("Score both on out-of-sample CVaR(0.95) of daily emissions.", None, False),
    ("Pre-registered: if spatial structure matters, the joint schedule must win. "
     "It does not.", NAVY, True),
], size=18, gap=14, lead="Shuffled marginals")
_pic(s, FIG / "schedule_us_west.png", 7.0, 1.6, max_w=5.7, max_h=4.6)
_text(s, 7.0, 6.3, 5.7, 0.4, [[(
    "Joint vs shuffled schedules, side by side.", SANS, 12, MUTED, False)]],
    align=PP_ALIGN.CENTER)

# ---- 18. Backup: tail dependence ----
s = content("Backup: non-elliptical tail dependence", 18)
_pic(s, FIG / "tail_dependence_taskc.png", 0.7, 1.55, max_w=7.2, max_h=4.9)
_bullets(s, 8.0, 1.7, 4.7, 4.8, [
    ("Residual dependence is non-elliptical: regions go clean together more than "
     "dirty together (chi_L > chi_U).", None, False),
    ("A covariance ball forces chi_L = chi_U by construction.", None, False),
    ("So an elliptical ambiguity set is blind to the only structure left, "
     "and it sits in the clean tail a risk-averse scheduler ignores.", NAVY, True),
], size=17, gap=14, lead="Structural reason covariance fails")

# =============================================================================== save
out = DECK / "capstone_defense.pptx"
prs.save(str(out))
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
